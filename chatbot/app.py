import streamlit as st
import pandas as pd
import os
import faiss
import numpy as np
from io import BytesIO
from langchain.text_splitter import CharacterTextSplitter
from langchain.docstore.in_memory import InMemoryDocstore
from dotenv import load_dotenv
from langchain_core.messages import HumanMessage, SystemMessage
from sentence_transformers import SentenceTransformer, util
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_huggingface import HuggingFaceEndpoint
from langchain.chains import RetrievalQA
from PyPDF2 import PdfReader
from docx import Document
from langchain_groq import ChatGroq
from pptx import Presentation  # For handling PPTX files

# Load environment variables from .env file
load_dotenv()
huggingface_api_key = os.getenv('HUGGINGFACEHUB_API_TOKEN')
groq_api_key = os.getenv('GROQ_API_KEY')

if huggingface_api_key is None:
    raise ValueError("Hugging Face API token not found. Please check your .env file for the correct key.")

# Set environment variable for Hugging Face
os.environ['HUGGINGFACEHUB_API_TOKEN'] = huggingface_api_key

# Initialize session state variables for chat history
if 'chat_history' not in st.session_state:
    st.session_state['chat_history'] = []

# App title
st.title("LLM RAG CHATBOT")

# Sidebar for chat history
st.sidebar.title("Chat History")
for i, (question, answer) in enumerate(st.session_state['chat_history']):
    if st.sidebar.button(f"Question {i+1}: {question}"):
        st.write("**Question**:", question)
        st.write("**Answer**:", answer)

# Function to process various input types
def process_input(input_type, input_data):
    # Code for processing PDF, PPTX, Text, DOCX, and TXT inputs
    loader = None
    if input_type == "PDF":
        pdf_reader = PdfReader(input_data)
        text = "".join([page.extract_text() for page in pdf_reader.pages])
        documents = text
    elif input_type == "PPTX":
        ppt = Presentation(input_data)
        text = ""
        for slide in ppt.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text += shape.text + "\n"
        documents = text
    elif input_type == "Text":
        documents = input_data
    elif input_type == "DOCX":
        doc = Document(input_data)
        text = "\n".join([para.text for para in doc.paragraphs])
        documents = text
    elif input_type == "TXT":
        text = input_data.read().decode('utf-8')
        documents = text
    else:
        raise ValueError("Unsupported input type")

    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=100)
    texts = text_splitter.split_text(documents)

    hf_embeddings = HuggingFaceEmbeddings(
        model_name="sentence-transformers/all-mpnet-base-v2",
        model_kwargs={'device': 'cpu'},
        encode_kwargs={'normalize_embeddings': False}
    )
    # Create FAISS index and vector store
    sample_embedding = np.array(hf_embeddings.embed_query("sample text"))
    dimension = sample_embedding.shape[0]
    index = faiss.IndexFlatL2(dimension)
    vector_store = FAISS(
        embedding_function=hf_embeddings.embed_query,
        index=index,
        docstore=InMemoryDocstore(),
        index_to_docstore_id={},
    )
    vector_store.add_texts(texts)
    return vector_store

# Function to find the most relevant rows based on the query
def find_relevant_rows(query, data, embeddings, model, top_k=5):
    query_embedding = model.encode([query], convert_to_tensor=True)
    cos_scores = util.pytorch_cos_sim(query_embedding, embeddings)[0]
    top_results = np.argpartition(-cos_scores, range(top_k))[0:top_k].cpu().numpy()
    return data.iloc[top_results]

# Function to answer a question based on the vector store or CSV data
def answer_question(input_type, vectorstore, query, data=None, embeddings=None, model=None):
    if input_type == "CSV":
        relevant_rows = find_relevant_rows(query, data, embeddings, model)
        context = relevant_rows.to_string(index=False)
        llm = ChatGroq(groq_api_key=groq_api_key, model_name="llama3-70b-8192", temperature=0.2)
        messages = [
            SystemMessage(content="You are a helpful assistant. Use the provided context to answer the user's question in a concise way."),
            HumanMessage(content=f"Context: {context}\n\nQuestion: {query}")
        ]
        result = llm(messages)
        return result.content
    else:
        llm = HuggingFaceEndpoint(repo_id='meta-llama/Meta-Llama-3-8B-Instruct', token=huggingface_api_key, temperature=0.6)
        qa = RetrievalQA.from_chain_type(llm=llm, retriever=vectorstore.as_retriever())
        answer = qa({"query": query})
        return answer['result']

# Function to load CSV data and compute embeddings
def load_data_and_embeddings(uploaded_file):
    data = pd.read_csv(uploaded_file)
    model = SentenceTransformer('distilbert-base-nli-stsb-mean-tokens')
    embeddings = model.encode(data.astype(str).values.tolist(), convert_to_tensor=True, show_progress_bar=True)
    return data, embeddings, model

# Main app logic
input_type = st.selectbox("Choose Input Type", ["PDF", "Text", "DOCX", "TXT", "CSV", "PPTX"])
input_data = None

if input_type == "CSV":
    uploaded_file = st.file_uploader("Upload your CSV file", type=["csv"])
    if uploaded_file:
        data, embeddings, model = load_data_and_embeddings(uploaded_file)
        st.info("CSV loaded successfully")
        st.dataframe(data.head(3), use_container_width=True)
else:
    input_data = st.file_uploader("Upload your file", type=['pdf', 'docx', 'txt', 'pptx']) if input_type in ["PDF", "DOCX", "TXT", "PPTX"] else st.text_input("Enter text")

if st.button("Proceed"):
    if input_type == "CSV":
        st.session_state["vectorstore"] = None  # Not using vector store for CSV
    else:
        st.session_state["vectorstore"] = process_input(input_type, input_data)

# User query and response handling
query = st.text_input("Ask your question")
if st.button("Submit"):
    if input_type == "CSV" and uploaded_file:
        answer = answer_question(input_type, None, query, data, embeddings, model)
    elif "vectorstore" in st.session_state and st.session_state["vectorstore"]:
        answer = answer_question(input_type, st.session_state["vectorstore"], query)
    else:
        st.error("Please upload a valid file and click Proceed.")
    
    # Display answer and add to chat history
    if answer:
        st.write("**Answer**:", answer)
        st.session_state['chat_history'].append((query, answer))
